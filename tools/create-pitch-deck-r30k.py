#!/usr/bin/env python3
"""
Create SiyaBusa ERP Investor Pitch Deck - R30,000 Ask
February 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Colors
DARK_BLUE = RGBColor(0x1a, 0x36, 0x5d)  # #1a365d
ACCENT_BLUE = RGBColor(0x38, 0x82, 0xf6)  # #3882f6
LIGHT_BLUE = RGBColor(0xeb, 0xf4, 0xff)  # #ebf4ff
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xe5, 0x3e, 0x3e)  # #e53e3e
GREEN = RGBColor(0x38, 0xa1, 0x69)  # #38a169
ORANGE = RGBColor(0xed, 0x8a, 0x36)  # #ed8a36

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, DARK_BLUE)
    
    # Main title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "SiyaBusa ERP"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Power. SME Pricing."
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Ask
    ask = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.6))
    tf = ask.text_frame
    p = tf.paragraphs[0]
    p.text = "Seeking R30,000 Seed Investment"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
    tf = date.text_frame
    p = tf.paragraphs[0]
    p.text = "February 2026 | Confidential"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

def add_problem_slide(prs):
    """Slide 2: The Problem - Emotional"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Painful Truth"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Emotional statement
    statement = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    tf = statement.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\"I wake up at 3am worried about my business.\nI know I need better systems, but I can't afford SAP.\nSo I stay stuck with Excel spreadsheets and hope nothing falls through the cracks.\""
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    attribution = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.4))
    tf = attribution.text_frame
    p = tf.paragraphs[0]
    p.text = "— Every SME owner in South Africa"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    
    # The numbers
    numbers = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.5))
    tf = numbers.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "95,000"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "South African SMEs are STUCK"
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    # Gap explanation
    gap = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(8.4), Inches(1.8))
    tf = gap.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌  QuickBooks = Too basic (just accounting, no real ERP)"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    
    p2 = tf.add_paragraph()
    p2.text = "❌  SAP/Oracle = Too expensive (R50,000+/month)"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = "❌  No middle ground = Businesses suffer in silence"
    p3.font.size = Pt(16)
    p3.font.color.rgb = DARK_GRAY
    p3.space_before = Pt(8)

def add_solution_slide(prs):
    """Slide 3: Our Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "We Built The Answer"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Product name
    product = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1))
    tf = product.text_frame
    p = tf.paragraphs[0]
    p.text = "SiyaBusa ERP"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Value prop
    value = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.6))
    tf = value.text_frame
    p = tf.paragraphs[0]
    p.text = "Full Enterprise ERP at SME Prices"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Stats boxes - Left
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.2), Inches(2.6), Inches(1.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0x2d, 0x4a, 0x7c)
    left_box.line.fill.background()
    tf = left_box.text_frame
    tf.paragraphs[0].text = "25+"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Modules"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Stats boxes - Center
    center_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(3.2), Inches(2.6), Inches(1.5))
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = RGBColor(0x2d, 0x4a, 0x7c)
    center_box.line.fill.background()
    tf = center_box.text_frame
    tf.paragraphs[0].text = "400+"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "API Endpoints"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Stats boxes - Right
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(3.2), Inches(2.6), Inches(1.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0x2d, 0x4a, 0x7c)
    right_box.line.fill.background()
    tf = right_box.text_frame
    tf.paragraphs[0].text = "✅ LIVE"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "In Production"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # Price comparison
    price = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.2))
    tf = price.text_frame
    p = tf.paragraphs[0]
    p.text = "R1,500/month vs SAP at R50,000+/month"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "siyabusaerp.co.za"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

def add_what_built_slide(prs):
    """Slide 4: What We've Built"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Complete ERP System — Already Built"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Left column - Core
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(2.8))
    tf = left.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Core Business Modules"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    modules_left = [
        "✅ Financial Accounting (GL, AP, AR)",
        "✅ Inventory Management",
        "✅ Sales & CRM",
        "✅ Purchase Management", 
        "✅ HR & Payroll (SARS-ready)",
        "✅ Manufacturing (BOM)"
    ]
    for m in modules_left:
        p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    # Right column - Verticals
    right = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(2.8))
    tf = right.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Industry Verticals"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    modules_right = [
        "✅ Healthcare",
        "✅ Construction",
        "✅ Logistics & Fleet",
        "✅ Agriculture",
        "✅ Mining",
        "✅ Property Management"
    ]
    for m in modules_right:
        p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    # Tech stack
    tech = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(2))
    tf = tech.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "Backend: Node.js + TypeScript  |  Frontend: React + Vite"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = "Database: PostgreSQL  |  AI: xAI Grok  |  Hosting: Vultr (Johannesburg)"
    p3.font.size = Pt(14)
    p3.font.color.rgb = DARK_GRAY
    p3.space_before = Pt(4)
    
    p4 = tf.add_paragraph()
    p4.text = "Email: Resend  |  SMS: Twilio  |  Video: Daily.co"
    p4.font.size = Pt(14)
    p4.font.color.rgb = DARK_GRAY
    p4.space_before = Pt(4)
    
    # Status
    status = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6), Inches(4), Inches(0.6))
    status.fill.solid()
    status.fill.fore_color.rgb = GREEN
    status.line.fill.background()
    tf = status.text_frame
    tf.paragraphs[0].text = "🚀 LIVE NOW: siyabusaerp.co.za"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_ask_slide(prs):
    """Slide 5: The Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Ask"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Amount
    amount = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    tf = amount.text_frame
    p = tf.paragraphs[0]
    p.text = "R30,000"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "for 3 Months → First Customer by March 2026"
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    # Budget breakdown
    budget = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(3.5))
    tf = budget.text_frame
    tf.word_wrap = True
    
    items = [
        ("🖥️  Infrastructure (Vultr)", "R6,000", "20%"),
        ("💻  Development (GitHub)", "R2,050", "7%"),
        ("📧  Communications (Resend + Twilio)", "R2,550", "8.5%"),
        ("🤖  AI Services (xAI Grok)", "R2,400", "8%"),
        ("📢  Marketing (LinkedIn + Google)", "R10,500", "35%"),
        ("🚗  Operations", "R6,500", "21.5%"),
    ]
    
    p = tf.paragraphs[0]
    p.text = f"{items[0][0]:<45} {items[0][1]:>10}    {items[0][2]}"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Courier New"
    
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = f"{item[0]:<45} {item[1]:>10}    {item[2]}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Courier New"
        p.space_before = Pt(6)
    
    # Total line
    p = tf.add_paragraph()
    p.text = "─" * 60
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = f"{'TOTAL':<45} {'R30,000':>10}    100%"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Courier New"

def add_timeline_slide(prs):
    """Slide 6: Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "3-Month Journey to Revenue"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Month boxes
    months = [
        ("FEBRUARY", "Launch", "Marketing campaigns go live\nFirst demos scheduled", ACCENT_BLUE),
        ("MARCH", "🎉 First Customer!", "Sign first paying client\nR1,500+ MRR", GREEN),
        ("APRIL", "Scale", "Second customer\nPipeline growing", ORANGE),
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    
    for i, (month, title_text, desc, color) in enumerate(months):
        # Month box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(1.5), Inches(3), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2d, 0x4a, 0x7c)
        box.line.fill.background()
        
        # Month name
        month_txt = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.7), Inches(3), Inches(0.5))
        tf = month_txt.text_frame
        p = tf.paragraphs[0]
        p.text = month
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_txt = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2.2), Inches(3), Inches(0.6))
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_txt = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(3), Inches(2.6), Inches(1.8))
        tf = desc_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        p.alignment = PP_ALIGN.CENTER
    
    # Bottom stats
    stats = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = stats.text_frame
    p = tf.paragraphs[0]
    p.text = "Target: R15,000 MRR by August 2026  |  Break-even: 7-10 customers"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_why_low_risk_slide(prs):
    """Slide 7: Why Low Risk"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Why This is Low Risk"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Risk items
    risks = [
        ("1", "Product is DONE", "Not a prototype. 25 modules, 400+ APIs, live in production."),
        ("2", "Infrastructure is CHEAP", "Vultr costs 70% less than AWS. R2,000/month total."),
        ("3", "Market is READY", "95,000 SMEs need affordable ERP. Post-COVID digitization."),
        ("4", "Amount is SMALL", "R30,000 for 3 months. Low risk, high potential."),
        ("5", "Timeline is SHORT", "60 days to first revenue. Quick validation."),
    ]
    
    y_pos = 1.3
    for num, title_text, desc in risks:
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y_pos), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Title and description
        txt = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos - 0.05), Inches(8), Inches(0.8))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY
        
        y_pos += 1.0

def add_pricing_slide(prs):
    """Slide 8: Pricing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Simple, Profitable Pricing"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Pricing tiers
    tiers = [
        ("Starter", "R1,500", "Up to 5 users", "Small businesses"),
        ("Professional", "R3,500", "Up to 15 users", "Growing companies"),
        ("Enterprise", "R7,500", "Unlimited users", "Large operations"),
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    colors = [ACCENT_BLUE, GREEN, ORANGE]
    
    for i, (name, price, users, target) in enumerate(tiers):
        # Tier box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(1.5), Inches(3), Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2d, 0x4a, 0x7c)
        box.line.fill.background()
        
        # Tier name
        name_txt = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.7), Inches(3), Inches(0.5))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        
        # Price
        price_txt = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2.3), Inches(3), Inches(0.8))
        tf = price_txt.text_frame
        p = tf.paragraphs[0]
        p.text = price
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "/month"
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p2.alignment = PP_ALIGN.CENTER
        
        # Users
        users_txt = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(3.5), Inches(3), Inches(0.8))
        tf = users_txt.text_frame
        p = tf.paragraphs[0]
        p.text = users
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = target
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p2.alignment = PP_ALIGN.CENTER
    
    # Comparison
    compare = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tf = compare.text_frame
    p = tf.paragraphs[0]
    p.text = "Compare to SAP: R50,000+/month"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "We're 90% cheaper with the same features"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = GREEN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

def add_infrastructure_slide(prs):
    """Slide 9: Infrastructure"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Infrastructure: Vultr Cloud"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Why Vultr
    why = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    tf = why.text_frame
    p = tf.paragraphs[0]
    p.text = "Why Vultr over AWS?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Benefits
    benefits = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.8))
    tf = benefits.text_frame
    tf.word_wrap = True
    
    items = [
        "✅ Johannesburg datacenter = Fast for SA users",
        "✅ 70% cheaper than AWS for same specs",
        "✅ No hidden data transfer fees",
        "✅ Simple pricing, easy to scale"
    ]
    
    p = tf.paragraphs[0]
    p.text = items[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
    
    # Server breakdown
    servers = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    tf = servers.text_frame
    p = tf.paragraphs[0]
    p.text = "Server Configuration (R2,000/month)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Server table
    table_txt = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(7), Inches(2))
    tf = table_txt.text_frame
    tf.word_wrap = True
    
    servers_data = [
        ("Main Production", "8GB / 4vCPU", "SiyaBusa ERP"),
        ("Database", "4GB / 2vCPU", "PostgreSQL"),
        ("Additional Apps", "4GB / 2vCPU", "Future projects"),
        ("Backup/Staging", "2GB / 1vCPU", "DR + testing"),
    ]
    
    for i, (name, spec, purpose) in enumerate(servers_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {name}: {spec} — {purpose}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        if i > 0:
            p.space_before = Pt(6)

def add_cta_slide(prs):
    """Slide 10: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Build Something Great"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # The ask again
    ask = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
    tf = ask.text_frame
    p = tf.paragraphs[0]
    p.text = "R30,000 → First Customer → March 2026"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Steps
    steps = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(2))
    tf = steps.text_frame
    tf.word_wrap = True
    
    step_items = [
        "1️⃣  See it live — 30-minute demo at siyabusaerp.co.za",
        "2️⃣  Review the numbers — Detailed budget attached",
        "3️⃣  Fund — R30,000 seed investment",
        "4️⃣  Celebrate — First client in March! 🎉"
    ]
    
    p = tf.paragraphs[0]
    p.text = step_items[0]
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    for item in step_items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)

def add_contact_slide(prs):
    """Slide 11: Contact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    
    # Logo/Name
    name = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = name.text_frame
    p = tf.paragraphs[0]
    p.text = "SiyaBusa ERP"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.6))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Power. SME Pricing."
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    tf = contact.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Sibusiso Buncele"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Founder & CEO"
    p2.font.size = Pt(18)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    
    p3 = tf.add_paragraph()
    p3.text = "Masaphokati Technologies (Pty) Ltd"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(16)
    
    p4 = tf.add_paragraph()
    p4.text = "🌐 siyabusaerp.co.za"
    p4.font.size = Pt(18)
    p4.font.color.rgb = ACCENT_BLUE
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(16)
    
    p5 = tf.add_paragraph()
    p5.text = "📧 info@masaphokati.co.za"
    p5.font.size = Pt(18)
    p5.font.color.rgb = DARK_GRAY
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(8)
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Confidential | February 2026"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    add_title_slide(prs)
    add_problem_slide(prs)
    add_solution_slide(prs)
    add_what_built_slide(prs)
    add_ask_slide(prs)
    add_timeline_slide(prs)
    add_why_low_risk_slide(prs)
    add_pricing_slide(prs)
    add_infrastructure_slide(prs)
    add_cta_slide(prs)
    add_contact_slide(prs)
    
    # Save
    output_path = "/workspaces/WorldClass-ERP/investor-docs/SiyaBusa-ERP-Investor-Pitch-R30K.pptx"
    prs.save(output_path)
    print(f"✅ Pitch deck created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
